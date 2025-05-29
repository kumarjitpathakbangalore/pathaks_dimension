import os
import json # Import json for saving summaries
from pptx import Presentation
import openai # Keep import here for type hinting/clarity even if client is dynamic
import matplotlib.pyplot as plt
from tqdm import tqdm
# Import specific client libraries, but only initialize them if selected
# from groq import Groq # Not needed at the top level if imported conditionally
# import google.generativeai as genai # Not needed at the top level if imported conditionally


class SlideSummarizer:
    """
    Extracts text from PowerPoint presentations (.pptx) in a specified folder,
    summarizes each slide using a chosen Language Model (LLM) provider (OpenAI, Groq, Gemini),
    and can optionally save these summaries.
    """

    def __init__(self, folder_path, model="gpt-4o", provider="openai", save_path=None):
        """
        Initializes the SlideSummarizer.

        Args:
            folder_path (str): The path to the folder containing .pptx files.
            model (str): The specific LLM model to use (e.g., "gpt-4o", "llama3-8b-8192", "gemini-pro").
                         Defaults to "gpt-4o".
            provider (str): The LLM service provider to use ('openai', 'groq', 'gemini').
                            Defaults to "openai".
            save_path (str, optional): The path to a JSON file where summaries will be saved.
                                       If None, summaries are not saved to a file.
        
        Raises:
            ValueError: If an unsupported LLM provider is specified.
            KeyError: If the required API key environment variable is not set.
        """
        self.folder_path = folder_path
        self.model = model
        self.provider = provider.lower()
        self.save_path = save_path
        self.client = None # Initialize client to None

        # --- Initialize LLM Client based on provider ---
        if self.provider == "openai":
            api_key = os.environ.get("OPENAI_API_KEY")
            if not api_key:
                raise KeyError("OPENAI_API_KEY environment variable not set.")
            self.client = openai.OpenAI(api_key=api_key) # Use openai.OpenAI() for new client
            print("🚀 Initialized OpenAI client.")
        elif self.provider == "groq":
            from groq import Groq # Import here to avoid unnecessary dependency if not used
            api_key = os.environ.get("GROQ_API_KEY")
            if not api_key:
                raise KeyError("GROQ_API_KEY environment variable not set.")
            self.client = Groq(api_key=api_key)
            print("🚀 Initialized Groq client.")
        elif self.provider == "gemini":
            import google.generativeai as genai # Import here
            api_key = os.environ.get("GOOGLE_API_KEY")
            if not api_key:
                raise KeyError("GOOGLE_API_KEY environment variable not set.")
            genai.configure(api_key=api_key)
            self.client = genai # The client is the genai module itself for Gemini
            print("🚀 Initialized Gemini client.")
        else:
            raise ValueError(
                f"Unsupported provider: '{self.provider}'. Choose from: 'openai', 'groq', 'gemini'."
            )

        # Ensure the folder path exists
        if not os.path.isdir(self.folder_path):
            raise FileNotFoundError(f"Folder not found: {self.folder_path}")
            
        # Create directory for saving summaries if save_path is provided
        if self.save_path:
            output_dir = os.path.dirname(self.save_path)
            if output_dir and not os.path.exists(output_dir):
                os.makedirs(output_dir, exist_ok=True)
                print(f"📁 Created directory for summaries: {output_dir}")

    def extract_pptx_text(self):
        """
        Extracts all textual content from each slide of .pptx files
        within the specified folder.

        Returns:
            dict: A dictionary where keys are PowerPoint filenames (e.g., "presentation.pptx")
                  and values are lists of strings, where each string is the concatenated
                  text content of a single slide.
        """
        extracted_texts = {}
        print(f"📖 Extracting text from PPTX files in: {self.folder_path}")
        for file in os.listdir(self.folder_path):
            if file.endswith(".pptx"):
                full_path = os.path.join(self.folder_path, file)
                try:
                    prs = Presentation(full_path)
                    slides_text = []
                    for i, slide in enumerate(prs.slides):
                        slide_content = []
                        for shape in slide.shapes:
                            # Check if the shape has text and if the text is not empty
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_content.append(shape.text.strip())
                        
                        # Join all text elements on a slide into a single string
                        text = " ".join(slide_content)
                        slides_text.append(text)
                    extracted_texts[file] = slides_text
                    print(f"✅ Extracted text from '{file}' ({len(slides_text)} slides).")
                except Exception as e:
                    print(f"❌ Error extracting text from '{file}': {e}")
        return extracted_texts

    def summarize_slide(self, slide_text):
        """
        Summarizes the given slide text using the configured LLM.
        Handles empty content and truncates long inputs for API limits.

        Args:
            slide_text (str): The raw text content of a single slide.

        Returns:
            str: A 2-line summary of the slide content, or an error message
                 if summarization fails or content is empty.
        """
        if not slide_text.strip():
            return "[Slide has no detectable content, consider its purpose visually.]"

        # Truncate long slide texts to avoid exceeding model token limits
        # A common limit is around 4096 tokens, 3000 chars is a safe buffer.
        if len(slide_text) > 3000:
            slide_text = slide_text[:3000] + " [Content truncated...]"

        # Craft the prompt for the LLM
        # Instructs for a 2-line summary, avoids repetition, and suggests analyzing empty slides.
        prompt = (
            "Summarize the following slide content in exactly two concise lines. "
            "Do not use phrases like '2-line summary' or 'summary of the slide'. "
            "Be direct and to the point. If the content is sparse, infer the slide's likely purpose. "
            "Also, anticipate potential questions a user might have about this slide's topic.\n\n"
            f"Slide Content:\n{slide_text}"
        )

        try:
            if self.provider == "openai" or self.provider == "groq":
                # Both OpenAI and Groq use a similar chat completions API
                response = self.client.chat.completions.create(
                    model=self.model,
                    messages=[{"role": "user", "content": prompt}],
                    temperature=0.3, # Lower temperature for more factual, less creative summaries
                )
                return response.choices[0].message.content.strip()

            elif self.provider == "gemini":
                # Gemini uses GenerativeModel for content generation
                model_instance = self.client.GenerativeModel(self.model)
                response = model_instance.generate_content(prompt)
                return response.text.strip()

        except Exception as e:
            # Catch any API errors or network issues during summarization
            return f"[Error summarizing slide: {e}]"

    def summarize_all(self):
        """
        Orchestrates the extraction of text from all PPTX files in the folder
        and summarizes each slide. Optionally saves the summaries to a JSON file.

        Returns:
            dict: A nested dictionary where keys are filenames and values are lists
                  of summarized strings for each slide in that file.
                  Example: {
                      "presentation1.pptx": ["Summary of slide 1", "Summary of slide 2"],
                      "presentation2.pptx": ["Summary of slide 1"]
                  }
        """
        extracted_texts = self.extract_pptx_text()
        all_summaries = {}

        if not extracted_texts:
            print("❗ No PPTX files found or no text extracted. No summaries to generate.")
            return all_summaries

        print("\n📝 Starting slide summarization for all extracted content...")
        for filename, slides_text_list in extracted_texts.items():
            print(f"Processing '{filename}' with {len(slides_text_list)} slides...")
            file_summaries = []
            # Use tqdm for a progress bar during summarization
            for slide_index, slide_text in enumerate(tqdm(slides_text_list, desc=f"Summarizing {filename}")):
                summary = self.summarize_slide(slide_text)
                file_summaries.append(summary)
            all_summaries[filename] = file_summaries
            print(f"✅ Finished summarizing '{filename}'.")

        if self.save_path:
            self._save_summaries(all_summaries)

        return all_summaries

    def _save_summaries(self, summaries):
        """
        Internal helper method to save the generated summaries to a JSON file.

        Args:
            summaries (dict): The dictionary of summaries to save.
        """
        try:
            with open(self.save_path, "w", encoding="utf-8") as f:
                json.dump(summaries, f, indent=4)
            print(f"💾 Successfully saved all summaries to: {self.save_path}")
        except IOError as e:
            print(f"❌ Error saving summaries to '{self.save_path}': {e}")
        except Exception as e:
            print(f"❌ An unexpected error occurred while saving summaries: {e}")

    def display_image(self, img):
        """
        Displays a PIL Image object using Matplotlib.
        Note: This method seems out of place in a class focused on text summarization.
              Consider moving it to a more appropriate display/visualization utility class.

        Args:
            img (PIL.Image.Image): The image to display.
        """
        plt.figure(figsize=(10, 7)) # Adjust figure size as needed
        plt.imshow(img)
        plt.axis('off') # Hide axes for cleaner image presentation
        plt.show()